

from pptx import Presentation
from pptx.util import Cm, Pt



with open('source.txt', 'r', encoding='UTF8') as file:    # hello.txt 파일을 읽기 모드(r)로 열기
    lines = file.readlines()


prs = Presentation()



left = Cm(4.23)
top = Cm(10.08)
width =Cm(3)
height = Cm(3)

layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

num=2
for i in lines:
    if i=="\n":
        layout = prs.slide_layouts[num]
        slide = prs.slides.add_slide(layout)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        num=num+1
    else:
        p = tf.add_paragraph()
        p.text = i.split("\n")[0]
        p.font.size = Pt(44)      




# 프레젠테이션 파일 저장
prs.save("./test.pptx")